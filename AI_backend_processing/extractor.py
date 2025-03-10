import pdfplumber
import docx
from pptx import Presentation

def extract_text_from_pdf(pdf_file):
    """Extract text from a PDF file received as an uploaded file object."""
    text = ""
    with pdfplumber.open(pdf_file) as pdf:
        for page in pdf.pages:
            extracted_text = page.extract_text()
            if extracted_text:  # Ensure text is not None
                text += extracted_text + "\n"
    return text.strip()

def extract_text_from_docx(docx_file):
    """Extract text from a DOCX file received as an uploaded file object."""
    doc = docx.Document(docx_file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(pptx_file):
    """Extract text from a PPTX file received as an uploaded file object."""
    text = ""
    prs = Presentation(pptx_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()
